from pptx import Presentation
from pptx.util import Inches
import os
from pdf_converter import convert_to_pdf
from send_pdf import send_email
    
def modify_pptx(type, template_path, output_path, replacements, image_path=None):
    ppt = Presentation(template_path)
    for slide in ppt.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for key, value in replacements.items():
                            if key in run.text:
                                run.text = run.text.replace(key, value)
            
            if shape.has_text_frame and "Picture" in shape.text:
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                
                sp = shape._element
                sp.getparent().remove(sp)
                
                slide.shapes.add_picture(image_path, left, top, width, height)
                
    ppt.save(output_path)

def create_pptx(data, template_path, output_path, image_file_path=None):
    replacements = {
        "Name": data["name"],
        "Position": data["position"],
        "Localization": data["localization"],
        "Description": data["description"],
    }
    
    modify_pptx(data["type"], template_path, output_path, replacements, image_file_path)
    pdf_output_path = output_path.replace(".pptx", ".pdf")
    convert_to_pdf(output_path, pdf_output_path)

    
    
